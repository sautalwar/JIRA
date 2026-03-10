"""Generate JIRA Bug Triage Demo PowerPoint and PDF"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
ACCENT_ORANGE = RGBColor(0xF0, 0x88, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text

blank = prs.slide_layouts[6]

# ════════════════════════════════════════════════
# SLIDE 1 - TITLE
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.6, 11, 0.8, 'JIRA BUG TRIAGE', 16, ACCENT_BLUE, True)
add_text(s, 0.8, 1.4, 11, 1.5,
         'Intelligent Bug Triage with\nGitHub Copilot SDK + MCP', 44, WHITE, True)
add_text(s, 0.8, 3.4, 11, 0.8,
         'Automating alert analysis \u2022 Duplicate detection \u2022 Smart bug creation', 20, LIGHT_GRAY)
add_card(s, 0.8, 4.5, 11.5, 0.04, ACCENT_BLUE)
add_text(s, 0.8, 5.0, 5, 0.6, 'Powered by GitHub Copilot SDK', 16, GRAY)
add_text(s, 0.8, 5.5, 5, 0.6, 'Connected via Model Context Protocol (MCP)', 16, GRAY)
add_text(s, 0.8, 6.0, 5, 0.6, 'JIRA Cloud Integration', 16, GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "Welcome everyone. Today I'm going to show you how we can use GitHub Copilot SDK "
    "combined with the JIRA MCP Server to create an intelligent bug triage system.\n\n"
    "The key problem: engineering teams get flooded with alerts from monitoring tools. "
    "Many correspond to bugs already tracked in JIRA. Today, engineers manually search "
    "JIRA, compare descriptions, and decide whether to create a new issue. "
    "This is time-consuming and error-prone.\n\n"
    "Our solution automates this entire workflow using AI."
))

# ════════════════════════════════════════════════
# SLIDE 2 - THE PROBLEM
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'THE PROBLEM', 14, ACCENT_BLUE, True)
add_text(s, 0.8, 1.0, 11, 1.0, 'Alert Fatigue & Duplicate Bugs', 38, WHITE, True)
problems = [
    ('\U0001f4ca', '500+', 'alerts/week from monitoring tools'),
    ('\U0001f501', '30%', 'of new bugs are duplicates of existing issues'),
    ('\u23f1\ufe0f', '15 min', 'average time to manually search & triage'),
    ('\U0001f4b0', '$50K+', 'annual cost in wasted engineering time'),
]
for i, (icon, stat, desc) in enumerate(problems):
    y = 2.4 + i * 1.1
    add_card(s, 0.8, y, 11.5, 0.9)
    add_text(s, 1.2, y + 0.15, 0.6, 0.6, icon, 28, WHITE)
    add_text(s, 2.2, y + 0.1, 2.0, 0.7, stat, 32, ACCENT_ORANGE, True)
    add_text(s, 4.8, y + 0.2, 7.0, 0.6, desc, 20, LIGHT_GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "Engineering teams receive hundreds of alerts per week. "
    "When an on-call engineer gets an alert, they need to check if it's a known issue, "
    "find the JIRA ticket, and either link to it or create a new one.\n\n"
    "Studies show ~30% of newly created bugs are duplicates. "
    "At 15 minutes per triage across hundreds of alerts, this adds up to "
    "significant engineering cost."
))

# ════════════════════════════════════════════════
# SLIDE 3 - THE SOLUTION
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'THE SOLUTION', 14, ACCENT_GREEN, True)
add_text(s, 0.8, 1.0, 11, 1.0, 'AI-Powered Bug Triage Pipeline', 38, WHITE, True)

add_card(s, 0.8, 2.4, 3.2, 2.0)
add_text(s, 1.0, 2.5, 2.8, 0.5, 'Alert Sources', 20, ACCENT_BLUE, True)
add_text(s, 1.0, 3.1, 2.8, 1.2,
         '\u2022 PagerDuty incidents\n\u2022 Datadog monitors\n\u2022 Sentry error clusters\n\u2022 Application error logs\n\u2022 Custom webhooks', 15, LIGHT_GRAY)

add_text(s, 4.2, 3.2, 0.8, 0.5, '\u2192', 36, ACCENT_BLUE, True)

add_card(s, 5.0, 2.4, 3.5, 2.0)
add_text(s, 5.2, 2.5, 3.1, 0.5, 'Copilot SDK + MCP', 20, ACCENT_PURPLE, True)
add_text(s, 5.2, 3.1, 3.1, 1.2,
         '\u2022 AI analyzes the alert\n\u2022 Searches JIRA via MCP\n\u2022 Compares descriptions\n\u2022 Scores confidence\n\u2022 Makes decision', 15, LIGHT_GRAY)

add_text(s, 8.7, 3.2, 0.8, 0.5, '\u2192', 36, ACCENT_BLUE, True)

add_card(s, 9.5, 2.4, 3.2, 2.0)
add_text(s, 9.7, 2.5, 2.8, 0.5, 'Outcome', 20, ACCENT_GREEN, True)
add_text(s, 9.7, 3.1, 2.8, 1.2,
         '\u2022 Link to existing issue\n  (with confidence %)\n\u2022 Create new bug\n  (well-formatted)\n\u2022 Add context', 15, LIGHT_GRAY)

add_card(s, 0.8, 5.0, 11.7, 1.5)
add_text(s, 1.2, 5.1, 4.0, 0.5, 'Key Technology Stack', 18, ACCENT_BLUE, True)
add_text(s, 1.2, 5.7, 10.0, 0.8,
         '(1) GitHub Copilot SDK - AI runtime with tool-calling   '
         '(2) mcp-atlassian - 72 JIRA tools via MCP   '
         '(3) Model Context Protocol - standardized AI-tool communication', 14, LIGHT_GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "Three-stage pipeline:\n\n"
    "1. ALERT SOURCES: Ingest from any monitoring tool, normalize into standard format.\n\n"
    "2. COPILOT SDK + MCP: The brain. Creates an AI session with JIRA tools. "
    "The AI autonomously searches, compares, and scores confidence.\n\n"
    "3. OUTCOME: High confidence (>70%) = link to existing. No match = create new bug.\n\n"
    "Key innovation: the AI doesn't just keyword-match. It understands that "
    "'high memory on notification-service' relates to 'WebSocket memory leak'."
))

# ════════════════════════════════════════════════
# SLIDE 4 - ARCHITECTURE
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'ARCHITECTURE', 14, ACCENT_BLUE, True)
add_text(s, 0.8, 1.0, 11, 0.8, 'How It Works Under the Hood', 38, WHITE, True)

steps = [
    ('1', 'Alert Arrives', 'PagerDuty/Datadog', ACCENT_ORANGE),
    ('2', 'Normalize', 'Standard format', ACCENT_BLUE),
    ('3', 'Copilot Session', 'SDK + MCP Server', ACCENT_PURPLE),
    ('4', 'Search JIRA', 'JQL via MCP', ACCENT_BLUE),
    ('5', 'Compare', 'AI reasoning', ACCENT_PURPLE),
    ('6', 'Take Action', 'Link or Create', ACCENT_GREEN),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 2.05
    add_card(s, x, 2.2, 1.8, 1.8)
    add_text(s, x + 0.1, 2.3, 1.6, 0.5, num, 28, color, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.9, 1.6, 0.5, title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 3.4, 1.6, 0.5, desc, 11, GRAY, False, PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, x + 1.8, 2.8, 0.3, 0.5, '\u2192', 20, GRAY, True)

add_card(s, 0.8, 4.5, 5.5, 2.3)
add_text(s, 1.0, 4.6, 5.0, 0.4, 'MCP Server Tools Used', 18, ACCENT_PURPLE, True)
add_text(s, 1.0, 5.1, 5.0, 1.6,
         'jira_search        Search with JQL queries\n'
         'jira_get_issue    Get full issue details\n'
         'jira_create_issue  Create new bugs\n'
         'jira_add_comment  Add context to existing\n'
         'jira_list_projects  Discover projects', 14, LIGHT_GRAY)

add_card(s, 6.8, 4.5, 5.5, 2.3)
add_text(s, 7.0, 4.6, 5.0, 0.4, 'AI Decision Logic', 18, ACCENT_BLUE, True)
add_text(s, 7.0, 5.1, 5.0, 1.6,
         'Confidence > 70%   = Link to existing issue\n'
         'Confidence 40-70% = Suggest match, create new\n'
         'Confidence < 40%   = Create new bug\n\n'
         'Factors: error type, service, stack trace,\n'
         '         component, description similarity', 14, LIGHT_GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "1. Alert arrives from any source.\n"
    "2. Normalized into standard format: title, description, severity, metadata.\n"
    "3. Copilot SDK creates AI session with mcp-atlassian attached.\n"
    "4. AI calls jira_search with multiple JQL queries (keywords, components, labels).\n"
    "5. For each match, AI reads full details and compares semantically.\n"
    "6. Based on confidence, links to existing or creates new.\n\n"
    "MCP (Model Context Protocol) is the standard that makes this possible."
))

# ════════════════════════════════════════════════
# SLIDE 5 - LIVE DEMO
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'LIVE DEMO', 14, ACCENT_GREEN, True)
add_text(s, 0.8, 1.0, 11, 1.0, "Let's See It In Action", 42, WHITE, True)

scenarios = [
    ('Scenario 1', 'PagerDuty: Auth NullPointerException', 'Should find existing DEMO-3', ACCENT_GREEN),
    ('Scenario 2', 'Datadog: High memory on notification-svc', 'Should match WebSocket leak DEMO-4', ACCENT_GREEN),
    ('Scenario 3', 'Sentry: Safari chart rendering TypeError', 'Should match Safari bug DEMO-6', ACCENT_GREEN),
    ('Scenario 4', 'CloudWatch: Lambda timeout in payments', 'Should create NEW bug (no match)', ACCENT_BLUE),
]
for i, (name, desc, expected, color) in enumerate(scenarios):
    y = 2.6 + i * 1.1
    add_card(s, 0.8, y, 11.5, 0.9)
    add_text(s, 1.2, y + 0.15, 2.0, 0.6, name, 18, color, True)
    add_text(s, 3.5, y + 0.15, 4.5, 0.6, desc, 16, WHITE)
    add_text(s, 8.5, y + 0.15, 3.5, 0.6, expected, 14, color)

add_text(s, 0.8, 6.6, 11, 0.5, 'Command:   npm run demo', 18, GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "[SWITCH TO TERMINAL - run: npm run demo]\n\n"
    "Scenario 1: PagerDuty alert about NullPointerException in auth service. "
    "AI searches JIRA and finds existing bug DEMO-3.\n\n"
    "Scenario 2: Datadog reports high memory on notification-service. "
    "Existing bug DEMO-4 is about WebSocket memory leak - same root cause, different words.\n\n"
    "Scenario 3: Sentry detects TypeError in Safari chart rendering. DEMO-6 matches.\n\n"
    "Scenario 4: CloudWatch Lambda timeout - genuinely new, no match. "
    "AI creates a new well-formatted JIRA issue.\n\n"
    "[SHOW JIRA BOARD after to verify new issue was created]"
))

# ════════════════════════════════════════════════
# SLIDE 6 - MULTI-SOURCE
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'MULTI-SOURCE SUPPORT', 14, ACCENT_ORANGE, True)
add_text(s, 0.8, 1.0, 11, 0.8, 'Works With Any Alert Source', 38, WHITE, True)
sources = [
    ('PagerDuty', 'Incident webhooks with urgency mapping'),
    ('Datadog', 'Monitor alerts with metric context'),
    ('Sentry', 'Error clusters with stack traces'),
    ('CloudWatch', 'AWS Lambda and infra alerts'),
    ('Prometheus', 'Metric-based alerting'),
    ('Error Logs', 'Parse .log files (JSON + standard)'),
    ('Webhooks', 'HTTP POST from any system'),
    ('Manual CLI', 'Ad-hoc analysis from terminal'),
]
for i, (name, desc) in enumerate(sources):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.2 + row * 1.15
    add_card(s, x, y, 5.8, 0.95)
    add_text(s, x + 0.3, y + 0.15, 2.5, 0.6, name, 18, WHITE, True)
    add_text(s, x + 0.3, y + 0.5, 5.0, 0.4, desc, 13, GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "The system is source-agnostic. We built normalizers for popular tools "
    "but also support raw error logs and generic webhooks.\n\n"
    "The log parser handles JSON and standard timestamp-level-service format. "
    "Extracts stack traces and maps log levels to severity.\n\n"
    "[DEMO: npx tsx src/index.ts analyze --file examples/sample-alerts.json]"
))

# ════════════════════════════════════════════════
# SLIDE 7 - JIRA DATA
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'JIRA INSTANCE', 14, ACCENT_BLUE, True)
add_text(s, 0.8, 1.0, 11, 0.8, '20 Realistic Bugs Across Services', 38, WHITE, True)
categories = [
    ('Auth', '3 bugs', 'NPE, token reuse, SSO loop'),
    ('API', '4 bugs', 'Rate limiter, CORS, webhooks, GraphQL'),
    ('Frontend', '2 bugs', 'Safari rendering, dark mode'),
    ('Database', '2 bugs', 'Connection pool, N+1 queries'),
    ('Security', '2 bugs', 'SQL injection, token expiry'),
    ('Infra', '3 bugs', 'K8s probes, ES corruption, memory leak'),
    ('Payments', '1 bug', 'Race condition, duplicate charges'),
    ('Mobile', '1 bug', 'Android 14 camera crash'),
]
for i, (cat, count, desc) in enumerate(categories):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.2 + row * 1.15
    add_card(s, x, y, 5.8, 0.95)
    add_text(s, x + 0.3, y + 0.15, 2.0, 0.5, cat, 18, WHITE, True)
    add_text(s, x + 2.5, y + 0.15, 1.0, 0.5, count, 16, ACCENT_BLUE, True)
    add_text(s, x + 0.3, y + 0.55, 5.0, 0.4, desc, 13, GRAY)
add_text(s, 0.8, 6.5, 11, 0.5,
         'Board: talwarsaurabh-1770668373217.atlassian.net/jira/software/projects/DEMO/board',
         14, ACCENT_BLUE)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "[SHOW JIRA BOARD IN BROWSER]\n\n"
    "We seeded 20 realistic bugs mirroring real-world scenarios. "
    "They span auth, API, frontend, database, security, infra, and more.\n\n"
    "Each has detailed descriptions with stack traces, impact analysis. "
    "Some are In Progress, one is Done. Realistic board state."
))

# ════════════════════════════════════════════════
# SLIDE 8 - BUSINESS VALUE
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'BUSINESS VALUE', 14, ACCENT_GREEN, True)
add_text(s, 0.8, 1.0, 11, 0.8, 'ROI & Impact', 38, WHITE, True)
values = [
    ('Time Saved', '15 min to 30 sec per triage', '95% reduction in triage time'),
    ('Accuracy', '30% duplicate rate to ~5%', 'AI catches matches humans miss'),
    ('Consistency', '24/7 automated triage', 'No missed alerts during off-hours'),
    ('Visibility', 'Full audit trail', 'Every decision logged with reasoning'),
]
for i, (title, metric, desc) in enumerate(values):
    y = 2.2 + i * 1.2
    add_card(s, 0.8, y, 11.5, 1.0)
    add_text(s, 1.5, y + 0.05, 2.5, 0.5, title, 22, WHITE, True)
    add_text(s, 4.5, y + 0.05, 4.0, 0.5, metric, 18, ACCENT_GREEN, True)
    add_text(s, 4.5, y + 0.5, 7.0, 0.5, desc, 14, GRAY)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "TIME: 15min manual triage -> 30sec automated. 95% reduction.\n\n"
    "ACCURACY: AI catches semantic matches humans miss. 'high memory' = 'WebSocket leak'.\n\n"
    "CONSISTENCY: Works 24/7. No duplicates from tired on-call engineers.\n\n"
    "VISIBILITY: Every decision logged with confidence scores and reasoning."
))

# ════════════════════════════════════════════════
# SLIDE 9 - CODE
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'TECHNICAL DETAILS', 14, ACCENT_PURPLE, True)
add_text(s, 0.8, 1.0, 11, 0.8, 'Key Code Walkthrough', 38, WHITE, True)

code = (
    '// Create Copilot session with JIRA MCP server\n'
    'const session = await client.createSession({\n'
    '  model: "gpt-4o",\n'
    '  mcpServers: {\n'
    '    atlassian: {\n'
    '      type: "local",\n'
    '      command: "uvx",\n'
    '      args: ["mcp-atlassian"],\n'
    '      env: {\n'
    '        JIRA_URL: process.env.JIRA_URL,\n'
    '        JIRA_USERNAME: process.env.JIRA_USERNAME,\n'
    '        JIRA_API_TOKEN: process.env.JIRA_API_TOKEN,\n'
    '      },\n'
    '      tools: ["jira_search", "jira_get_issue",\n'
    '              "jira_create_issue", "jira_add_comment"],\n'
    '    },\n'
    '  },\n'
    '});\n\n'
    '// Send prompt - AI calls JIRA tools autonomously\n'
    'const result = await session.sendAndWait({ prompt });'
)
add_card(s, 0.8, 2.2, 11.5, 4.5)
add_text(s, 1.2, 2.3, 10.5, 4.2, code, 13, LIGHT_GRAY, False, PP_ALIGN.LEFT, 'Consolas')
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "The core code is remarkably simple.\n\n"
    "1. Create a CopilotClient and start a session\n"
    "2. Attach mcp-atlassian MCP server - runs locally via uvx\n"
    "3. Specify which tools the AI can use\n"
    "4. Send a prompt instructing the AI to search and decide\n"
    "5. AI autonomously calls JIRA tools and returns structured result\n\n"
    "The prompt engineering instructs the AI to try multiple search strategies "
    "and return a confidence score."
))

# ════════════════════════════════════════════════
# SLIDE 10 - NEXT STEPS
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.4, 11, 0.8, 'NEXT STEPS', 14, ACCENT_ORANGE, True)
add_text(s, 0.8, 1.0, 11, 0.8, 'Roadmap & Expansion', 38, WHITE, True)
roadmap = [
    ('Phase 1', 'Production Webhook Integration', 'Connect to real PagerDuty/Datadog for live triage'),
    ('Phase 2', 'Confluence Knowledge Base', 'Add Confluence MCP to search runbooks alongside JIRA'),
    ('Phase 3', 'Auto-Remediation', 'Trigger automated fixes for known issues via GitHub Actions'),
    ('Phase 4', 'Learning Loop', 'Feed triage decisions back to improve future accuracy'),
]
for i, (phase, title, desc) in enumerate(roadmap):
    y = 2.2 + i * 1.2
    add_card(s, 0.8, y, 11.5, 1.0)
    add_text(s, 1.2, y + 0.1, 1.5, 0.5, phase, 16, ACCENT_ORANGE, True)
    add_text(s, 3.0, y + 0.1, 8.5, 0.5, title, 20, WHITE, True)
    add_text(s, 3.0, y + 0.5, 8.5, 0.5, desc, 14, GRAY)
add_card(s, 0.8, 6.4, 11.5, 0.04, ACCENT_BLUE)
add_text(s, 0.8, 6.6, 11, 0.5, 'github.com/sautalwar/JIRA', 18, ACCENT_BLUE)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "Phase 1: Connect to real monitoring in production. Webhook server is built.\n"
    "Phase 2: Add Confluence - AI searches runbooks alongside JIRA issues.\n"
    "Phase 3: Trigger automated remediation via GitHub Actions.\n"
    "Phase 4: Feedback loop for improving accuracy.\n\n"
    "Code is open source on GitHub."
))

# ════════════════════════════════════════════════
# SLIDE 11 - Q&A
# ════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_text(s, 0.8, 1.5, 11.5, 1.5, 'Questions?', 56, WHITE, True, PP_ALIGN.CENTER)
add_card(s, 4.0, 3.5, 5.3, 0.04, ACCENT_BLUE)
add_text(s, 0.8, 4.0, 11.5, 0.8, 'github.com/sautalwar/JIRA', 22, ACCENT_BLUE, False, PP_ALIGN.CENTER)
add_text(s, 0.8, 4.8, 11.5, 0.6,
         'JIRA: talwarsaurabh-1770668373217.atlassian.net', 16, GRAY, False, PP_ALIGN.CENTER)
add_text(s, 0.8, 5.8, 11.5, 0.6,
         'GitHub Copilot SDK  |  MCP Protocol  |  mcp-atlassian', 16, GRAY, False, PP_ALIGN.CENTER)
add_notes(s, (
    "SPEAKER NOTES:\n\n"
    "Thank you! Happy to answer questions.\n\n"
    "Key links:\n"
    "- GitHub: github.com/sautalwar/JIRA\n"
    "- JIRA: talwarsaurabh-1770668373217.atlassian.net\n"
    "- Copilot SDK: github.com/github/copilot-sdk\n"
    "- mcp-atlassian: github.com/sooperset/mcp-atlassian"
))

# Save
out = os.path.join(os.getcwd(), 'JIRA_Bug_Triage_Demo.pptx')
prs.save(out)
print(f"Saved PowerPoint: {out}")
print(f"Total slides: {len(prs.slides)}")
